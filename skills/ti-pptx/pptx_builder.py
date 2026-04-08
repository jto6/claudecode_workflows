"""
TI Presentation Builder
Core module for creating TI-branded PowerPoint presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
from typing import List, Tuple, Optional, Dict, Any
from dataclasses import dataclass
import os


# ──────────────────────────────────────────────────────────────
# Official TI Brand Color Palette
# Source: TI Color and Design Guide (2024-08-05)
# ──────────────────────────────────────────────────────────────

class TIColors:
    """Official TI corporate color palette.

    Usage rules (from TI Brand Guide):
      - Embrace red — RED / DARK_RED can fill large areas and headers.
      - RED + WHITE together for high-contrast, large-area layouts.
      - Teal family used sparingly — "a little teal goes a long way."
      - Lines and borders are encouraged — consistent with TI design style.

    IMPORTANT: Do not use colors outside this palette in TI-branded
    presentations. If you need a color not listed here, ask the user.
    """

    # Primary colors
    RED         = RGBColor(0xCC, 0x00, 0x00)  # TI Red (PMS 3546C) — primary brand
    BLACK       = RGBColor(0x00, 0x00, 0x00)  # Process Black
    GRAY        = RGBColor(0xAA, 0xAA, 0xAA)  # TI Gray (PMS Cool Gray 8C)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)  # Process White

    # Secondary colors
    DARK_RED    = RGBColor(0x99, 0x00, 0x00)  # PMS 1807 — dark accent, headers on dark bg
    DARK_TEAL   = RGBColor(0x11, 0x55, 0x66)  # PMS 5473 — section headers, dark cards
    TEAL        = RGBColor(0x00, 0x7C, 0x8C)  # PMS 321 — accent, use sparingly
    BRIGHT_CYAN = RGBColor(0x00, 0xBB, 0xCC)  # PMS 3215 — small highlights only
    LIGHT_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)  # PMS 5455 — card backgrounds, borders

    # Derived / convenience aliases
    BODY_TEXT   = BLACK                        # Default body text color
    CARD_BG     = LIGHT_GRAY                   # Light card/box background
    DARK_CARD   = DARK_TEAL                    # Dark card/box background


class TIFonts:
    """TI corporate font sizing for slide content.

    These sizes match the TI template typography hierarchy.
    Use these instead of raw Pt() values.
    """

    TITLE       = Pt(22)   # Slide titles
    SUBTITLE    = Pt(16)   # Subtitles
    SECTION     = Pt(16)   # Section headers within a slide
    BODY        = Pt(12)   # Body text
    SMALL_BODY  = Pt(11)   # Sub-bullets, secondary body
    SMALL       = Pt(10)   # Captions, card labels
    LABEL       = Pt(9)    # Fine print, annotations
    FAMILY      = None     # None = inherit from template slide master


# Legacy aliases for backward compatibility
TI_RED = TIColors.RED
TI_BLUE = TIColors.DARK_TEAL  # Corrected: was #00338D, now official dark teal
TI_DARK_GRAY = TIColors.BLACK
TI_LIGHT_GRAY = TIColors.GRAY

# Template locations - use bundled templates
SKILL_DIR = Path(__file__).parent
TEMPLATE_DIR = SKILL_DIR / "templates"

# Validate that templates directory exists and contains templates
if not TEMPLATE_DIR.exists():
    raise FileNotFoundError(
        f"Templates directory not found at {TEMPLATE_DIR}. "
        "The templates directory must be bundled with this skill."
    )

if not list(TEMPLATE_DIR.glob("*.pptx")):
    raise FileNotFoundError(
        f"No PowerPoint templates found in {TEMPLATE_DIR}. "
        "Please ensure template files (*.pptx) are present in the templates directory."
    )

# Template mappings
TEMPLATES = {
    "nda": "Presentation1_NDA_Restrictions.pptx",
    "internal": "Presentation2_NDA_Restrictions_InternalOnly.pptx",
    "selective": "Presentation3_SelectiveDisclosure.pptx",
    "max": "Presentation4_MaxRestrictions.pptx",
    "max_internal": "Presentation5_MaxRestrictions_InternalOnly.pptx"
}


@dataclass
class ColorScheme:
    """Custom color scheme for presentations."""
    primary: RGBColor = None
    accent: RGBColor = None
    text_primary: RGBColor = None
    text_secondary: RGBColor = None

    def __post_init__(self):
        if self.primary is None:
            self.primary = TIColors.RED
        if self.accent is None:
            self.accent = TIColors.DARK_TEAL
        if self.text_primary is None:
            self.text_primary = TIColors.BLACK
        if self.text_secondary is None:
            self.text_secondary = TIColors.GRAY


class TIPresentationBuilder:
    """
    Builder class for creating TI-branded PowerPoint presentations.

    Example:
        builder = TIPresentationBuilder(template_type="nda")
        builder.add_title_slide("My Presentation", "Subtitle")
        builder.add_content_slide("Overview", [("Point 1", 0), ("Point 2", 0)])
        builder.save()
    """

    def __init__(
        self,
        template_type: str = "nda",
        custom_template: Optional[str] = None,
        output_path: str = "output/presentation.pptx",
        color_scheme: Optional[ColorScheme] = None
    ):
        """
        Initialize the presentation builder.

        Args:
            template_type: One of "nda", "max", "internal", "selective"
            custom_template: Optional path to custom .pptx template
            output_path: Where to save the final presentation
            color_scheme: Optional custom color scheme
        """
        self.template_type = template_type
        self.custom_template = custom_template
        self.output_path = output_path
        self.color_scheme = color_scheme or ColorScheme()

        # Load template
        self.prs = self._load_template()
        self.template_name = self._get_template_name()

        # Metadata
        self.metadata = {}

    def _load_template(self) -> Presentation:
        """Load the PowerPoint template."""
        # Use custom template if provided
        if self.custom_template and Path(self.custom_template).exists():
            return Presentation(self.custom_template)

        # Use standard TI template
        template_file = TEMPLATES.get(self.template_type, TEMPLATES["nda"])
        template_path = Path(TEMPLATE_DIR) / template_file

        if template_path.exists():
            return Presentation(str(template_path))
        else:
            # Fall back to blank presentation
            print(f"Warning: Template not found at {template_path}, using blank")
            return Presentation()

    def _get_template_name(self) -> str:
        """Get the name of the loaded template."""
        if self.custom_template:
            return Path(self.custom_template).name
        return TEMPLATES.get(self.template_type, "blank")

    @property
    def slide_count(self) -> int:
        """Return the number of slides in the presentation."""
        return len(self.prs.slides)

    def set_metadata(
        self,
        title: str,
        subtitle: Optional[str] = None,
        author: Optional[str] = None,
        date: Optional[str] = None
    ):
        """
        Set presentation metadata.

        Args:
            title: Presentation title
            subtitle: Optional subtitle
            author: Author name
            date: Date string
        """
        self.metadata = {
            'title': title,
            'subtitle': subtitle,
            'author': author,
            'date': date
        }

        # Set core properties
        self.prs.core_properties.title = title
        if author:
            self.prs.core_properties.author = author

    def add_title_slide(self, title: str, subtitle: str = ""):
        """
        Add a title slide.

        Args:
            title: Main title
            subtitle: Subtitle text
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle

    def add_content_slide(
        self,
        title: str,
        bullets: List[Tuple[str, int]] = None,
        highlight_text: Optional[str] = None,
        highlight_color: str = "blue",
        font_size: int = 14
    ):
        """
        Add a content slide with bullet points.

        Args:
            title: Slide title
            bullets: List of (text, indent_level) tuples
            highlight_text: Optional text to highlight at bottom
            highlight_color: Color for highlight ("blue" or "red")
            font_size: Base font size for bullets
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[4])  # Title and Content

        title_shape = slide.shapes.title
        title_shape.text = title

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()  # Clear default text

        # Add bullets
        if bullets:
            for i, (text, level) in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(font_size if level == 0 else font_size - 2)

        # Add highlight text
        if highlight_text:
            tf.add_paragraph()
            p = tf.add_paragraph()
            p.text = highlight_text
            p.font.bold = True
            p.font.size = Pt(font_size + 2)

            color = self.color_scheme.accent if highlight_color == "blue" else self.color_scheme.primary
            p.font.color.rgb = color

    def add_two_column_slide(
        self,
        title: str,
        left_bullets: List[Tuple] = None,
        right_image: Optional[str] = None,
        right_text: Optional[str] = None
    ):
        """
        Add a two-column slide (text left, image/text right).

        Args:
            title: Slide title
            left_bullets: List of tuples: (text, level) or (text, level, color, bold)
            right_image: Path to image for right column
            right_text: Alternative text for right column
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Two Content

        title_shape = slide.shapes.title
        title_shape.text = title

        # Left column: bullets
        left = slide.placeholders[1]
        tf = left.text_frame
        tf.clear()

        if left_bullets:
            for i, item in enumerate(left_bullets):
                # Unpack with defaults
                if len(item) == 2:
                    text, level = item
                    color, bold = None, False
                elif len(item) == 4:
                    text, level, color, bold = item
                else:
                    text, level = item[0], item[1]
                    color, bold = None, False

                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = text
                p.level = level
                p.font.size = Pt(13 if level == 0 else 11)

                if color == "blue":
                    p.font.color.rgb = self.color_scheme.accent
                elif color == "red":
                    p.font.color.rgb = self.color_scheme.primary

                if bold:
                    p.font.bold = True

        # Right column: image or text
        right = slide.placeholders[2]

        if right_image and Path(right_image).exists():
            # Clear placeholder and add image
            sp = right._element
            sp.getparent().remove(sp)

            # Add image to slide
            left_pos = right.left
            top_pos = right.top
            height = right.height
            slide.shapes.add_picture(right_image, left_pos, top_pos, height=height)

        elif right_text:
            right.text = right_text

    def add_diagram_slide(
        self,
        title: str,
        diagram_type: str,
        content: str,
        font_size: int = 11
    ):
        """
        Add a slide with an ASCII diagram or flowchart.

        Args:
            title: Slide title
            diagram_type: "ascii" or "mermaid"
            content: The diagram content
            font_size: Font size for the diagram
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[4])  # Title and Content

        title_shape = slide.shapes.title
        title_shape.text = title

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = content

        # Format as monospace
        for paragraph in tf.paragraphs:
            paragraph.font.name = "Courier New"
            paragraph.font.size = Pt(font_size)

    def add_image_slide(
        self,
        title: str,
        image_path: str,
        caption: Optional[str] = None
    ):
        """
        Add a slide with a large image.

        Args:
            title: Slide title
            image_path: Path to image file
            caption: Optional caption below image
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[7])  # Title Only

        title_shape = slide.shapes.title
        title_shape.text = title

        if Path(image_path).exists():
            # Add image centered
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)

            slide.shapes.add_picture(image_path, left, top, width=width)

            # Add caption if provided
            if caption:
                textbox = slide.shapes.add_textbox(
                    Inches(1), Inches(6.5), Inches(8), Inches(0.5)
                )
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = caption
                p.font.size = Pt(10)
                p.font.color.rgb = TI_LIGHT_GRAY
                p.font.italic = True
                p.alignment = PP_ALIGN.CENTER

    # ──────────────────────────────────────────────────────────
    # Brand constants exposed as instance attributes
    # ──────────────────────────────────────────────────────────

    colors = TIColors
    fonts = TIFonts

    # ──────────────────────────────────────────────────────────
    # Freeform slide support
    # ──────────────────────────────────────────────────────────

    def add_freeform_slide(self, title: Optional[str] = None,
                           layout_name: str = "Title Only"):
        """Return a raw slide for custom spatial layout.

        The slide inherits TI branding from the template slide master
        (footer, logo, background). The caller places shapes freely
        using add_branded_box(), add_accent_line(), add_text_box(),
        or raw python-pptx calls.

        ALL colors must come from builder.colors.* and ALL font sizes
        from builder.fonts.*. Do not use raw RGBColor() or Pt() values.

        Args:
            title: Optional slide title text.
            layout_name: Template layout name (default "Title Only").

        Returns:
            The python-pptx Slide object.
        """
        layout = self._get_layout_by_name(layout_name)
        slide = self.prs.slides.add_slide(layout)
        if title:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 0:  # title placeholder
                    shape.text = title
                    break
        return slide

    def _get_layout_by_name(self, name: str):
        """Find a slide layout by name, falling back to index 7."""
        for layout in self.prs.slide_layouts:
            if layout.name == name:
                return layout
        # Fallback: "Title Only" is typically index 7
        return self.prs.slide_layouts[7]

    def add_branded_box(self, slide, left, top, width, height,
                        fill_color=None, text=None, font_size=None,
                        font_color=None, bold=False, word_wrap=True,
                        alignment=PP_ALIGN.LEFT, vertical="top",
                        corner_radius=0.05):
        """Add a rounded rectangle with TI brand-compliant styling.

        Args:
            slide: The slide object (from add_freeform_slide).
            left, top, width, height: Position/size in Emu or Inches.
            fill_color: A TIColors value (default: TIColors.CARD_BG).
            text: Optional text content.
            font_size: A TIFonts value (default: TIFonts.BODY).
            font_color: A TIColors value (default: TIColors.BODY_TEXT).
            bold: Whether text is bold.
            word_wrap: Enable word wrapping (default True).
            alignment: PP_ALIGN value (default LEFT).
            vertical: Vertical anchor — "top", "middle", or "bottom".
            corner_radius: Corner rounding 0.0–0.5 (default 0.05).

        Returns:
            The shape object (for further customization).
        """
        fill_color = fill_color or TIColors.CARD_BG
        font_color = font_color or TIColors.BODY_TEXT
        font_size = font_size or TIFonts.BODY

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        shape.adjustments[0] = corner_radius

        if text is not None:
            tf = shape.text_frame
            tf.word_wrap = word_wrap
            tf.margin_left = Pt(8)
            tf.margin_right = Pt(8)
            tf.margin_top = Pt(6)
            tf.margin_bottom = Pt(6)

            anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}
            tf.vertical_anchor = anchor_map.get(vertical, MSO_ANCHOR.TOP)

            p = tf.paragraphs[0]
            p.alignment = alignment
            run = p.add_run()
            run.text = text
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = font_color

        return shape

    def add_accent_line(self, slide, left, top, width,
                        color=None, thickness=Emu(18288)):
        """Add a horizontal accent line (defaults to TI Red).

        Args:
            slide: The slide object.
            left, top, width: Position and width in Emu or Inches.
            color: A TIColors value (default: TIColors.RED).
            thickness: Line thickness (default ~2px).

        Returns:
            The shape object.
        """
        color = color or TIColors.RED
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, thickness
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text_box(self, slide, left, top, width, height,
                     text="", font_size=None, font_color=None,
                     bold=False, alignment=PP_ALIGN.LEFT,
                     word_wrap=True):
        """Add a text box with TI brand-compliant styling.

        Args:
            slide: The slide object.
            left, top, width, height: Position/size.
            text: Text content.
            font_size: A TIFonts value (default: TIFonts.BODY).
            font_color: A TIColors value (default: TIColors.BODY_TEXT).
            bold: Whether text is bold.
            alignment: PP_ALIGN value.
            word_wrap: Enable word wrapping.

        Returns:
            The text frame object (for adding paragraphs).
        """
        font_size = font_size or TIFonts.BODY
        font_color = font_color or TIColors.BODY_TEXT

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = font_color
        return tf

    def add_paragraph(self, text_frame, text, font_size=None,
                      font_color=None, bold=False,
                      space_before=Pt(0), space_after=Pt(0)):
        """Add a paragraph to an existing text frame.

        Use this to build multi-line content in branded boxes or
        text boxes created by add_branded_box() or add_text_box().

        Args:
            text_frame: The text frame to append to.
            text: Paragraph text.
            font_size: A TIFonts value (default: TIFonts.BODY).
            font_color: A TIColors value (default: TIColors.BODY_TEXT).
            bold: Whether text is bold.
            space_before: Space above paragraph.
            space_after: Space below paragraph.

        Returns:
            The paragraph object.
        """
        font_size = font_size or TIFonts.BODY
        font_color = font_color or TIColors.BODY_TEXT

        p = text_frame.add_paragraph()
        p.space_before = space_before
        p.space_after = space_after
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = font_color
        return p

    def add_images_from_directory(
        self,
        directory: str,
        pattern: str = "*.png",
        captions: bool = True
    ):
        """
        Add multiple images from a directory, one per slide.

        Args:
            directory: Directory containing images
            pattern: File pattern to match (e.g., "*.png")
            captions: Whether to use filename as caption
        """
        dir_path = Path(directory)
        if not dir_path.exists():
            print(f"Warning: Directory not found: {directory}")
            return

        images = sorted(dir_path.glob(pattern))

        for img_path in images:
            caption = img_path.stem.replace('_', ' ').title() if captions else None
            self.add_image_slide(
                title=caption or "Image",
                image_path=str(img_path),
                caption=caption
            )

    def create_from_outline(self, outline: Dict[str, Any]):
        """
        Create presentation from a structured outline dictionary.

        Args:
            outline: Dictionary with title, subtitle, and slides list

        Example:
            outline = {
                "title": "My Presentation",
                "subtitle": "Subtitle",
                "author": "Team",
                "slides": [
                    {"type": "title", "content": {"title": "...", "subtitle": "..."}},
                    {"type": "content", "content": {"title": "...", "bullets": [...]}}
                ]
            }
        """
        # Set metadata
        if 'title' in outline:
            self.set_metadata(
                title=outline.get('title'),
                subtitle=outline.get('subtitle'),
                author=outline.get('author'),
                date=outline.get('date')
            )

        # Add slides
        for slide_def in outline.get('slides', []):
            slide_type = slide_def.get('type')
            content = slide_def.get('content', {})

            if slide_type == 'title':
                self.add_title_slide(
                    title=content.get('title', ''),
                    subtitle=content.get('subtitle', '')
                )

            elif slide_type == 'content':
                self.add_content_slide(
                    title=content.get('title', ''),
                    bullets=content.get('bullets', []),
                    highlight_text=content.get('highlight_text'),
                    highlight_color=content.get('highlight_color', 'blue')
                )

            elif slide_type == 'two_column':
                self.add_two_column_slide(
                    title=content.get('title', ''),
                    left_bullets=content.get('left', []),
                    right_image=content.get('right_image')
                )

            elif slide_type == 'diagram':
                self.add_diagram_slide(
                    title=content.get('title', ''),
                    diagram_type=content.get('diagram_type', 'ascii'),
                    content=content.get('content', '')
                )

            elif slide_type == 'image':
                self.add_image_slide(
                    title=content.get('title', ''),
                    image_path=content.get('image_path', ''),
                    caption=content.get('caption')
                )

    def save(self) -> str:
        """
        Save the presentation to file.

        Returns:
            Path to saved file
        """
        # Ensure output directory exists
        output_path = Path(self.output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Save presentation
        self.prs.save(str(output_path))

        return str(output_path)


# Utility function
def create_ti_presentation(
    title: str,
    slides: List[Dict[str, Any]],
    template_type: str = "nda",
    output_path: str = "output/presentation.pptx"
) -> str:
    """
    Quick helper to create a TI presentation.

    Args:
        title: Presentation title
        slides: List of slide definitions
        template_type: Template to use
        output_path: Output file path

    Returns:
        Path to created presentation
    """
    builder = TIPresentationBuilder(
        template_type=template_type,
        output_path=output_path
    )

    outline = {
        'title': title,
        'slides': slides
    }

    builder.create_from_outline(outline)
    return builder.save()
